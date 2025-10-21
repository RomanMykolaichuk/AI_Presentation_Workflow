"""Генератор базового PowerPoint-шаблону з плейсхолдерами у форматі {{name}}.

Приклад використання
--------------------
>>> from pathlib import Path
>>> from create_default_template import build_template
>>> build_template(Path("default_template.pptx"))
"""
from __future__ import annotations

from pathlib import Path

try:
    from pptx import Presentation
except ModuleNotFoundError as exc:  # pragma: no cover - потребує зовнішнього пакета
    raise ModuleNotFoundError(
        "Для генерації PPTX необхідно встановити python-pptx: pip install python-pptx"
    ) from exc


def build_template(target: Path) -> Path:
    """Створити презентацію з трьома слайдами та плейсхолдерами.

    Parameters
    ----------
    target: Path
        Шлях до файлу PPTX, який буде створено або перезаписано.

    Returns
    -------
    Path
        Шлях до створеного файлу.
    """

    prs = Presentation()

    # Слайд 1 — титул
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "{{title}}"
    slide.placeholders[1].text = "{{subtitle}}"

    # Слайд 2 — контент + зображення
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "{{section_title}}"
    body_shape = slide.shapes.placeholders[1]
    body_shape.text_frame.text = "{{body}}"
    picture_placeholder = slide.placeholders[1]
    if picture_placeholder.has_text_frame:
        picture_placeholder.text = "{{image:key_trends}}"

    # Слайд 3 — дві колонки
    layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(layout)
    left_title = slide.shapes.placeholders[0]
    left_body = slide.shapes.placeholders[1]
    right_title = slide.shapes.placeholders[2]
    right_body = slide.shapes.placeholders[3]
    left_title.text = "{{left_title}}"
    left_body.text = "{{left_body}}"
    right_title.text = "{{right_title}}"
    right_body.text = "{{right_body}}"

    image_shape = slide.shapes.add_textbox(left_body.left, right_body.top, right_body.width, right_body.height // 2)
    image_shape.text_frame.text = "{{image:comparison}}"

    prs.save(target)
    return target


if __name__ == "__main__":
    build_template(Path(__file__).with_name("default_template.pptx"))
    print("Шаблон згенеровано: default_template.pptx")
