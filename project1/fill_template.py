"""Модуль заповнення PowerPoint шаблону на основі структури.

Приклад використання
--------------------
>>> from fill_template import build_pptx
>>> build_pptx()
PosixPath('output/Presentation_demo.pptx')
"""
from __future__ import annotations

import json
import os
import tempfile
from pathlib import Path
from typing import Dict, Iterable, Optional

from dotenv import load_dotenv

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ModuleNotFoundError as exc:  # pragma: no cover - залежить від python-pptx
    raise ModuleNotFoundError(
        "Необхідно встановити python-pptx перед використанням fill_template: pip install python-pptx"
    ) from exc

try:
    from PIL import Image, ImageDraw, ImageFont
except ModuleNotFoundError as exc:  # pragma: no cover - залежить від pillow
    raise ModuleNotFoundError(
        "Необхідно встановити Pillow перед використанням fill_template: pip install Pillow"
    ) from exc

import requests


class ImageProvider:
    """Базовий генератор зображень."""

    def generate(self, prompt: str, size: tuple[int, int]) -> Path:  # pragma: no cover - інтерфейс
        raise NotImplementedError


class GeminiImageProvider(ImageProvider):
    def __init__(self, api_key: Optional[str]) -> None:
        self.api_key = api_key
        self.endpoint = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-pro-exp:generateImage"

    def generate(self, prompt: str, size: tuple[int, int]) -> Path:
        if not self.api_key:
            raise RuntimeError("GEMINI_API_KEY не налаштовано для генерації зображень")
        headers = {"Authorization": f"Bearer {self.api_key}"}
        payload = {"prompt": prompt, "size": {"width": size[0], "height": size[1]}}
        response = requests.post(self.endpoint, headers=headers, json=payload, timeout=120)
        response.raise_for_status()
        data = response.json()
        image_bytes = None
        if "image" in data:
            blob = data["image"]
            if isinstance(blob, str):
                import base64

                image_bytes = base64.b64decode(blob)
            elif isinstance(blob, (bytes, bytearray)):
                image_bytes = bytes(blob)
        if not image_bytes:
            raise RuntimeError("Gemini не повернув дані зображення")
        tmp = Path(tempfile.mkstemp(suffix=".png")[1])
        tmp.write_bytes(image_bytes)
        return tmp


class OllamaImageProvider(ImageProvider):
    def __init__(self, model: str = "llava") -> None:
        self.model = model
        self.endpoint = "http://localhost:11434/api/generate"

    def generate(self, prompt: str, size: tuple[int, int]) -> Path:
        payload = {"model": self.model, "prompt": prompt, "size": size, "format": "b64"}
        response = requests.post(self.endpoint, json=payload, timeout=120)
        response.raise_for_status()
        data = response.json()
        blob = data.get("image") or data.get("data")
        if not blob:
            raise RuntimeError("Ollama не повернув зображення")
        import base64

        image_bytes = base64.b64decode(blob)
        tmp = Path(tempfile.mkstemp(suffix=".png")[1])
        tmp.write_bytes(image_bytes)
        return tmp


class PlaceholderImageProvider(ImageProvider):
    """Формує заглушку у разі відсутності доступу до моделей."""

    def generate(self, prompt: str, size: tuple[int, int]) -> Path:
        width, height = size
        width = max(width, 800)
        height = max(height, 600)
        image = Image.new("RGB", (width, height), color=(245, 245, 245))
        draw = ImageDraw.Draw(image)
        text = f"Placeholder\n{prompt[:80]}"
        try:
            font = ImageFont.load_default()
        except OSError:  # pragma: no cover - fallback
            font = None
        draw.multiline_text((20, 20), text, fill=(80, 80, 80), font=font, spacing=12)
        tmp = Path(tempfile.mkstemp(suffix=".png")[1])
        image.save(tmp)
        return tmp


def _select_image_provider() -> ImageProvider:
    load_dotenv()
    choice = os.getenv("IMAGE_GENERATION_PROVIDER", "gemini").lower()
    if choice == "ollama":
        return OllamaImageProvider(os.getenv("OLLAMA_MODEL", "llava"))
    try:
        return GeminiImageProvider(os.getenv("GEMINI_API_KEY"))
    except Exception:  # noqa: BLE001
        return PlaceholderImageProvider()


class TemplateRenderer:
    def __init__(
        self,
        presentation: Presentation,
        mapping: Dict[str, Dict[str, Dict[str, str]]],
        meta: Dict[str, object],
        image_provider: ImageProvider,
    ) -> None:
        self.presentation = presentation
        self.mapping = mapping
        self.meta = meta
        self.image_provider = image_provider

    def render(self, slides: Iterable[Dict[str, object]]) -> None:
        slide_lookup = {int(slide["num"]): slide for slide in slides}
        for index, ppt_slide in enumerate(self.presentation.slides, start=1):
            slide_data = slide_lookup.get(index)
            if not slide_data:
                continue
            placeholder_map = self.mapping.get("slides", {}).get(str(index), {})
            for shape in ppt_slide.shapes:
                if not shape.has_text_frame and shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                text = shape.text if shape.has_text_frame else ""
                if shape.has_text_frame and "{{image:" in text:
                    key = text.strip()
                    config = placeholder_map.get(key)
                    if not config:
                        continue
                    prompt = self._transform_value(slide_data, config) or "Створити ілюстрацію"
                    size = self._resolve_size(index, key)
                    image_path = self._safe_generate_image(prompt, size)
                    self._replace_with_image(ppt_slide, shape, image_path)
                elif text and text.startswith("{{") and text.endswith("}}"):  # текстовий плейсхолдер
                    config = placeholder_map.get(text)
                    if not config:
                        continue
                    rendered = self._transform_value(slide_data, config)
                    self._fill_text(shape, rendered)

    def _transform_value(self, slide: Dict[str, object], config: Dict[str, str]) -> str:
        field = config.get("field")
        value = slide.get(field, "") if field else ""
        transform = config.get("transform")
        if not transform:
            return str(value)
        if transform == "summary":
            return str(value).split("\n")[0][:400]
        if transform == "bullets":
            return str(value)
        if transform.startswith("split_title"):
            parts = str(slide.get("title", "")).split(":", 1)
            if transform.endswith("left"):
                return parts[0]
            return parts[1].strip() if len(parts) > 1 else parts[0]
        if transform.startswith("split_body"):
            paragraphs = str(slide.get("text", "")).split("\n")
            midpoint = max(1, len(paragraphs) // 2)
            if transform.endswith("left"):
                return "\n".join(paragraphs[:midpoint])
            return "\n".join(paragraphs[midpoint:])
        if transform == "comparison":
            return f"Порівняльна ілюстрація: {value}"
        return str(value)

    def _fill_text(self, shape, text: str) -> None:
        text_frame = shape.text_frame
        text_frame.clear()
        paragraphs = [line.strip() for line in text.split("\n") if line.strip()]
        if not paragraphs:
            paragraphs = [""]
        for idx, paragraph in enumerate(paragraphs):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = paragraph
            if len(paragraphs) > 1:
                p.level = 0

    def _resolve_size(self, slide_num: int, placeholder: str) -> tuple[int, int]:
        slides_meta = self.meta.get("slides", [])
        if 0 <= slide_num - 1 < len(slides_meta):
            candidate = slides_meta[slide_num - 1].get("placeholders", {})
            data = candidate.get(placeholder)
            if data and data.get("size_px"):
                return tuple(data["size_px"])
        for slide_meta in slides_meta:
            placeholder_data = slide_meta.get("placeholders", {})
            data = placeholder_data.get(placeholder)
            if data and data.get("size_px"):
                return tuple(data["size_px"])
        return (1024, 768)

    def _safe_generate_image(self, prompt: str, size: tuple[int, int]) -> Path:
        try:
            return self.image_provider.generate(prompt, size)
        except Exception:  # noqa: BLE001
            return PlaceholderImageProvider().generate(prompt, size)

    def _replace_with_image(self, slide, shape, image_path: Path) -> None:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        if shape.has_text_frame:
            shape.text_frame.clear()
        try:
            image_path.unlink()
        except OSError:
            pass


def build_pptx(
    structure_path: Path | str = Path("output/structure.json"),
    template_path: Path | str = Path("templates/default_template.pptx"),
    mapping_path: Path | str = Path("templates/default_template_mapping.json"),
    meta_path: Path | str = Path("templates/default_template.json"),
    output_dir: Path | str = Path("output"),
) -> Path:
    """Заповнити PPTX-шаблон згідно зі структурою та зберегти у output."""

    load_dotenv()
    structure_file = Path(structure_path)
    if not structure_file.exists():
        raise FileNotFoundError("Не знайдено output/structure.json. Спочатку виконайте generate_structure.generate")
    if not Path(template_path).exists():
        raise FileNotFoundError(
            "Не знайдено шаблон PPTX. Запустіть templates/create_default_template.py для його створення."
        )

    slides_payload = json.loads(structure_file.read_text(encoding="utf-8"))
    mapping = json.loads(Path(mapping_path).read_text(encoding="utf-8"))
    meta = json.loads(Path(meta_path).read_text(encoding="utf-8"))

    presentation = Presentation(str(template_path))
    renderer = TemplateRenderer(presentation, mapping, meta, _select_image_provider())
    renderer.render(slides_payload.get("slides", []))

    output_directory = Path(output_dir)
    output_directory.mkdir(parents=True, exist_ok=True)
    output_file = output_directory / f"Presentation_{os.getenv('PRESENTATION_SUFFIX', 'demo')}.pptx"
    presentation.save(str(output_file))
    return output_file


__all__ = ["build_pptx"]
